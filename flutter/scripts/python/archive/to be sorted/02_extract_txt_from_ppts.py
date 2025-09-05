#!/usr/bin/env python3
import sys, os
from pptx import Presentation
from projects.ramayanam_sriramghanapatigal import folders
from pathlib import Path
def get_filepaths(directory, file_type, skip_subdir=True):
    """Collect file paths of a certain type in a directory."""
    file_paths = []
    if skip_subdir:
        for filename in sorted(os.listdir(directory)):
            if not filename.startswith('.') and not filename.startswith('~') and filename.endswith(file_type):
                filepath = os.path.join(directory, filename)
                file_paths.append(filepath)
    else:
        for root, directories, files in os.walk(directory):
            for filename in files:
                if filename.endswith(file_type):
                    filepath = os.path.join(root, filename)
                    file_paths.append(filepath)
    return sorted(file_paths)

def get_metadata(input_file):
    with open(input_file, 'r', encoding="utf-8") as file:
        lines = file.readlines()
        meta_lines = []
        meta_begin = False
        for line in lines:
            line = line.strip()
            if line.startswith("-- METADATA") or line.startswith("--METADATA"):
                meta_begin = True
            elif line.startswith("--END METADATA") or line.startswith("-- END METADATA"):
                meta_begin = False
            meta_lines.append(line)
            if not meta_begin:
                break
    file.close()
    return meta_lines

def insert_bt(textbox):
    updated_text = ""
    i = 0
    para_count = len(textbox.paragraphs)
    while i < len(textbox.paragraphs) - 1:
        line = textbox.paragraphs[i].text
        next_line = textbox.paragraphs[i + 1].text
        if next_line.startswith('\t'):
            if para_count > 2:
                updated_text += line.strip() + "{b}{t}" + next_line.strip() + "\n"
            else:
                updated_text += line.strip() + "\n\t" + next_line.strip() + "\n"
            i += 2
        else:
            updated_text += line.strip() + "\n"
            i += 1
    return updated_text

def export_text_from_ppts(fldr_ppts, fldr_metadata, output_fldr_content, lang, shape_ids):

    # First export all the ppt files to text files with the same name
    ppt_full_paths = get_filepaths(fldr_ppts, ".pptm")

    for ppt_file_path in ppt_full_paths:
        basename = os.path.basename(ppt_file_path)
        basename_without_extension = os.path.splitext(basename)[0]

        # First write text files for videos
        input_txt_file = os.path.join(fldr_metadata, basename_without_extension + ".txt")
        meta_data = get_metadata(input_txt_file)

        Path(output_fldr_content).mkdir(parents=True, exist_ok=True)
        output_txt_file = os.path.join(output_fldr_content, basename_without_extension + ".txt")
        with (open(output_txt_file, 'w', encoding="utf-8") as file):
            for meta_line in meta_data:
                file.write(meta_line + "\n")
            prs = Presentation(ppt_file_path)
            slidecount = len(prs.slides)
            for i in range(1, slidecount - 1):
                slide = prs.slides[i]
                textbox = slide.shapes[shape_ids[lang]].text_frame  # change this to the desired textbox index
                slide_text = textbox.text + "\n"
                if lang == "ta":
                    textbox_sa =  slide.shapes[shape_ids["sa"]].text_frame
                    if not len(textbox_sa.paragraphs) == 4:
                        slide_text = insert_bt(textbox)

                file.write(slide_text)

        file.close()

if __name__ == "__main__":
    volume = 5
    shape_ids = {
        "sa": 0,
        "ta": 1,
    }

    fldr_ppts = folders[volume].FLDR_PPTS
    fldr_metadata = folders[volume].FLDR_METADATA

    output_fldr_content = folders[volume].STAGING_FLDR_CONTENT

    # EXPORT SANSKRIT
    # Let us export sanskrit and tamil at the same time
    export_text_from_ppts(fldr_ppts, fldr_metadata, output_fldr_content, "sa", shape_ids)

    #EXPORT TAMIL
    folders[volume].lang = "ta"
    fldr_ppts = folders[volume].FLDR_PPTS
    fldr_metadata = folders[volume].FLDR_METADATA
    output_fldr_content = folders[volume].STAGING_FLDR_CONTENT
    export_text_from_ppts(fldr_ppts, fldr_metadata, output_fldr_content,"ta",  shape_ids)

