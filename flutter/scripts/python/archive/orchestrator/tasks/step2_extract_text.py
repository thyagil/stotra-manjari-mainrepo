# orchestrator/tasks/step2_extract_text.py
import os
import subprocess
from pathlib import Path
from pptx import Presentation


def get_metadata(input_file):
    with open(input_file, 'r', encoding="utf-8") as file:
        lines = file.readlines()
        meta_lines = []
        meta_begin = False
        for line in lines:
            line = line.strip()
            if line.startswith("-- METADATA") or line.startswith("--METADATA"):
                meta_begin = True
            elif line.startswith("--END METADATA") or line.startswith("-- END METADATA"):
                meta_begin = False
            meta_lines.append(line)
            if not meta_begin:
                break
    return meta_lines


def convert_key_to_pptx(key_file):
    """
    Handle Keynote .key file conversion.
    Rules:
      - If .pptm exists → use it
      - If .pptx exists → use it
      - If neither exists → convert .key → .pptx
    Returns Path to the file that should be processed.
    """
    key_file = Path(key_file).resolve()
    pptx_file = key_file.with_suffix(".pptx")
    pptm_file = key_file.with_suffix(".pptm")

    if pptm_file.exists():
        print(f"⏭️ Found existing PPTM, using: {pptm_file}")
        return pptm_file

    if pptx_file.exists():
        print(f"⏭️ Found existing PPTX, using: {pptx_file}")
        return pptx_file

    # If neither exists, convert .key → .pptx
    print(f"🔄 Converting {key_file} → {pptx_file}")
    applescript = """
    on run argv
        set keyFile to POSIX file (item 1 of argv)
        set outFile to POSIX file (item 2 of argv)
        tell application "Keynote"
            open keyFile
            export front document to outFile as Microsoft PowerPoint
            close front document saving no
        end tell
    end run
    """
    subprocess.run(
        ["osascript", "-e", applescript, str(key_file), str(pptx_file)],
        check=True
    )
    return pptx_file


def run(args, resolver):
    """
    Step 2: Extract text from PPTM/PPTX/KEY files into .txt files.
    """
    for lang, shape_id in args.shape_id.items():
        export_text_from_ppts(args, resolver, lang, shape_id)


def export_text_from_ppts(args, resolver, lang, shape_id):
    input_dir = resolver.ppts_input_working
    files_to_process = []

    if args.project_type == "standalone":
        project_file = resolver.get_file_by_name(input_dir, args.project_name, [".pptm", ".pptx", ".key"])
        if project_file:
            files_to_process = [project_file]   # ensure list
    else:
        files_to_process = sorted(
            f.resolve()
            for f in Path(resolver.ppts_input_working).glob("*")
            if f.suffix.lower() in [".pptm", ".pptx", ".key"]
            and not f.name.startswith(('.', '~'))
        )

    for ppt_file_path in files_to_process:
        # Handle .key conversion
        if ppt_file_path.suffix.lower() == ".key":
            ppt_file_path = convert_key_to_pptx(ppt_file_path)

        basename_without_extension = Path(ppt_file_path).stem

        # metadata
        meta_data = None
        if args.include_metadata_text:
            input_txt_file = Path(resolver.metadata_input_working) / f"{basename_without_extension}.txt"
            if input_txt_file.exists():
                meta_data = get_metadata(input_txt_file)

        # output
        output_fldr_content = Path(resolver.content_output_working) / lang
        output_fldr_content.mkdir(parents=True, exist_ok=True)

        output_txt_file = output_fldr_content / f"{basename_without_extension}.txt"
        with open(output_txt_file, "w", encoding="utf-8") as file:
            if meta_data:
                for meta_line in meta_data:
                    file.write(meta_line + "\n")

            prs = Presentation(ppt_file_path)
            for i in range(1, len(prs.slides)):
                slide = prs.slides[i]

                # make sure shape_id is valid
                if shape_id >= len(slide.shapes):
                    print(f"⚠️ Slide {i}: shape_id {shape_id} not found → stopping processing")
                    break

                shape = slide.shapes[shape_id]

                # check if shape has a text frame
                if not shape.has_text_frame:
                    print(f"⚠️ Slide {i}: shape {shape_id} has no text frame → stopping processing")
                    break

                textbox = shape.text_frame
                if textbox:
                    slide_text = textbox.text + "\n"
                    file.write(slide_text)

        print(f"📖 Extracted {output_txt_file}")
