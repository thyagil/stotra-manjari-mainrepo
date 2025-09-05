import re
import subprocess
from pathlib import Path
from pptx import Presentation
import plistlib
from typing import Optional
from admin.config.utils import get_project_path


def get_metadata(input_file: Path):
    with open(input_file, 'r', encoding="utf-8") as file:
        lines = file.readlines()
        meta_lines = []
        meta_begin = False
        for line in lines:
            line = line.strip()
            if line.startswith("-- METADATA") or line.startswith("--METADATA"):
                meta_begin = True
            elif line.startswith("--END METADATA") or line.startswith("-- END METADATA"):
                meta_begin = False
            meta_lines.append(line)
            if not meta_begin:
                break
    return meta_lines


def convert_key_to_pptx(key_file: Path) -> Path:
    key_file = key_file.resolve()
    pptx_file = key_file.with_suffix(".pptx")
    pptm_file = key_file.with_suffix(".pptm")

    if pptm_file.exists():
        print(f"⏭️ Found existing PPTM, using: {pptm_file}")
        return pptm_file
    if pptx_file.exists():
        print(f"⏭️ Found existing PPTX, using: {pptx_file}")
        return pptx_file

    print(f"🔄 Converting {key_file} → {pptx_file}")
    applescript = """
    on run argv
        set keyFile to POSIX file (item 1 of argv)
        set outFile to POSIX file (item 2 of argv)
        tell application "Keynote"
            open keyFile
            export front document to outFile as Microsoft PowerPoint
            close front document saving no
        end tell
    end run
    """
    subprocess.run(
        ["osascript", "-e", applescript, str(key_file), str(pptx_file)],
        check=True
    )
    return pptx_file


def extract_chapter_num(filename: str, regex: str) -> Optional[int]:
    """Extract chapter number using regex (supports multiple groups)."""
    match = re.search(regex, filename)
    if match:
        for g in match.groups():
            if g:
                try:
                    return int(g)
                except ValueError:
                    return None
    return None


def detect_language(text: str) -> Optional[str]:
    """Detect whether text is Sanskrit (Devanagari) or Tamil."""
    for ch in text:
        code = ord(ch)
        if 0x0900 <= code <= 0x097F:  # Devanagari
            return "sa"
        if 0x0B80 <= code <= 0x0BFF:  # Tamil
            return "ta"
    return None


def extract_ppt(
        ppt_file: Path,
        output_base: Path,
        project_format: str,
        chapter_regex: str,
        meta_prefix: str,
        meta_dirs: dict
):
    """Scan PPT, auto-detect Sanskrit/Tamil shapes, and export text into lang folders."""
    # ✅ Get chapter number consistently (standalone or regex)
    if project_format == "standalone":
        chap_num = 1
    else:
        chap_num = extract_chapter_num(ppt_file.stem, chapter_regex)

    if chap_num is None:
        print(f"⚠️ Could not parse chapter number from {ppt_file.name}")
        return

    chapter_id = f"chapter{chap_num:03d}"

    # ✅ Prepare output files for both langs
    lang_files = {}
    for lang in ["sa", "ta"]:
        out_dir = output_base / lang
        out_dir.mkdir(parents=True, exist_ok=True)
        out_file = out_dir / f"{chapter_id}.txt"
        lang_files[lang] = open(out_file, "w", encoding="utf-8")

        # ✅ Metadata lookup
        meta_data = None
        meta_dir = Path(meta_dirs.get(lang, "")) if meta_dirs.get(lang) else None
        if meta_dir:
            if meta_prefix:
                meta_file = meta_dir / f"{meta_prefix}{chap_num:02d}.txt"
                if meta_file.exists():
                    print(f"   📝 Found metadata file: {meta_file.name}")
                    meta_data = get_metadata(meta_file)
            if not meta_data:
                fallback_file = meta_dir / f"{chapter_id}.txt"
                if fallback_file.exists():
                    print(f"   📝 Found fallback metadata file: {fallback_file.name}")
                    meta_data = get_metadata(fallback_file)
        if meta_data:
            for meta_line in meta_data:
                lang_files[lang].write(meta_line + "\n")

    # ✅ Scan slides & shapes
    prs = Presentation(ppt_file)
    for i in range(1, len(prs.slides)-1):  # skip first slide
        slide = prs.slides[i]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            lang = detect_language(text)
            if lang in lang_files:
                lang_files[lang].write(text + "\n")
            else:
                print(f"⚠️ Slide {i}: Could not detect language → {text[:30]}...")

    for f in lang_files.values():
        f.close()

    print(f"📖 Extracted {ppt_file.name} → {output_base}")


def extract_text(project_id: str):
    """Step 2: Extract text for all volumes as per orchestrator.plist."""
    project_dir = get_project_path(project_id)
    orch_plist = project_dir / "orchestrator.plist"

    with open(orch_plist, "rb") as f:
        config = plistlib.load(f)

    project_format = config.get("format", "volume")
    staging_root = Path(config.get("staging_root", ""))
    if not staging_root:
        raise ValueError("❌ staging_root not defined in orchestrator.plist")

    step_cfg = config.get("text_extraction", {})
    volumes_cfg = step_cfg.get("volumes", {})
    step_regex = step_cfg.get("chapter_regex", r"(\d+)")  # ✅ step-level regex

    for vol_id, vol_cfg in volumes_cfg.items():
        if not vol_cfg.get("orchestrate", False):
            print(f"⏭️ Skipping {vol_id}")
            continue

        ppts_path = Path(vol_cfg.get("ppts_dir", ""))
        if not ppts_path.exists():
            print(f"⚠️ PPTs not found for {vol_id}: {ppts_path}")
            continue

        # Output path = staging_root/content/<volume>
        output_dir = staging_root / "content" / vol_id
        output_dir.mkdir(parents=True, exist_ok=True)

        files_to_process = [ppts_path] if ppts_path.is_file() else \
            sorted([
                f for f in ppts_path.glob("*")
                if f.suffix.lower() in [".pptm", ".pptx", ".key"]
                   and not f.name.startswith(('.', '~'))
            ])

        if not files_to_process:
            print(f"⚠️ No PPT/Keynote files in {ppts_path}")
            continue

        meta_prefix = vol_cfg.get("meta_prefix", "")
        meta_dirs = {
            "sa": vol_cfg.get("meta_dir_sa", ""),
            "ta": vol_cfg.get("meta_dir_ta", "")
        }

        for ppt_file in files_to_process:
            if ppt_file.suffix.lower() == ".key":
                ppt_file = convert_key_to_pptx(ppt_file)

            regex = vol_cfg.get("chapter_regex", step_regex)
            extract_ppt(ppt_file, output_dir, project_format, regex, meta_prefix, meta_dirs)


if __name__ == "__main__":
    extract_text("skp_srimad_bhagavatam")
