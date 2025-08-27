#!/usr/bin/env python3
import sys, os
from pptx import Presentation
from projects.ramayanam_sriramghanapatigal import folders


def get_filepaths(directory, file_type, skip_subdir=True):
    """Collect file paths of a certain type in a directory."""
    file_paths = []
    if skip_subdir:
        for filename in sorted(os.listdir(directory)):
            if not filename.startswith('.') and not filename.startswith('~') and filename.endswith(file_type):
                filepath = os.path.join(directory, filename)
                file_paths.append(filepath)
    else:
        for root, directories, files in os.walk(directory):
            for filename in files:
                if filename.endswith(file_type):
                    filepath = os.path.join(root, filename)
                    file_paths.append(filepath)
    return sorted(file_paths)

def export_text_from_ppts(ppt_folder, video_txt_folder, export_txt_folder_video, shape_id=0):
    """Export text from PPT slides into txt files."""
    os.makedirs(export_txt_folder_video, mode=0o777, exist_ok=True)

    ppt_full_paths = get_filepaths(ppt_folder, ".pptm")

    for ppt_file_path in ppt_full_paths:
        basename = os.path.basename(ppt_file_path)
        basename_without_extension = os.path.splitext(basename)[0]

        # Input metadata (if exists in video_txt_folder)
        input_txt_file = os.path.join(video_txt_folder, basename_without_extension + ".txt")
        metadata_lines = []
        if os.path.exists(input_txt_file):
            with open(input_txt_file, "r", encoding="utf-8") as f:
                metadata_lines = [line.strip() for line in f.readlines() if line.strip()]

        # Output text file
        output_txt_file = os.path.join(export_txt_folder_video, basename_without_extension + ".txt")
        with open(output_txt_file, 'w', encoding="utf-8") as file:
            # First write metadata
            for meta_line in metadata_lines:
                file.write(meta_line + "\n")

            prs = Presentation(ppt_file_path)
            slidecount = len(prs.slides)
            for i in range(1, slidecount - 1):  # skip first & last slide
                slide = prs.slides[i]
                textbox = slide.shapes[shape_id].text_frame  # shape index
                slide_text = textbox.text.strip() + "\n"
                file.write(slide_text)

        print(f"✅ Exported {output_txt_file}")


if __name__ == "__main__":

    volume = 5
    ppt_folder = folders[volume].FLDR_PPTS
    video_txt_folder = folders[volume].FLDR_CONTENT_VIDEO
    export_txt_folder_video = folders[volume].FLDR_CONTENT_FINAL

    # Let us export sanskrit and tamil at the same time
    export_text_from_ppts(ppt_folder, video_txt_folder, export_txt_folder_video, 0)

    folders[volume].lang = "ta"
    ppt_folder = folders[volume].FLDR_PPTS
    video_txt_folder = folders[volume].FLDR_CONTENT_VIDEO
    export_txt_folder_video = folders[volume].FLDR_CONTENT_FINAL
    export_text_from_ppts(ppt_folder, video_txt_folder, export_txt_folder_video, 1)

