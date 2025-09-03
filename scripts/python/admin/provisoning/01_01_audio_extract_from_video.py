import re
import subprocess
from pathlib import Path
from pptx import Presentation
import plistlib
from typing import Optional
from admin.config.utils import get_project_path


def get_metadata(input_file: Path):
    with open(input_file, 'r', encoding="utf-8") as file:
        lines = file.readlines()
    meta_lines, meta_begin = [], False
    for line in lines:
        line = line.strip()
        if line.startswith("-- METADATA") or line.startswith("--METADATA"):
            meta_begin = True
        elif line.startswith("--END METADATA") or line.startswith("-- END METADATA"):
            meta_begin = False
        meta_lines.append(line)
        if not meta_begin:
            break
    return meta_lines


def convert_key_to_pptx(key_file: Path) -> Path:
    key_file = key_file.resolve()
    pptx_file = key_file.with_suffix(".pptx")
    pptm_file = key_file.with_suffix(".pptm")
    if pptm_file.exists():
        return pptm_file
    if pptx_file.exists():
        return pptx_file
    applescript = """
    on run argv
        set keyFile to POSIX file (item 1 of argv)
        set outFile to POSIX file (item 2 of argv)
        tell application "Keynote"
            open keyFile
            export front document to outFile as Microsoft PowerPoint
            close front document saving no
        end tell
    end run
    """
    subprocess.run(["osascript", "-e", applescript, str(key_file), str(pptx_file)], check=True)
    return pptx_file


def extract_chapter_num(filename: str, regex: str) -> Optional[int]:
    """Extract chapter number using regex (supports multiple groups)."""
    match = re.search(regex, filename)
    if match:
        for g in match.groups():
            if g:
                try:
                    return int(g)
                except ValueError:
                    return None
    return None


def detect_language(text: str) -> Optional[str]:
    """Detect whether text is Sanskrit (Devanagari) or Tamil."""
    for ch in text:
        code = ord(ch)
        if 0x0900 <= code <= 0x097F:  # Devanagari
            return "sa"
        if 0x0B80 <= code <= 0x0BFF:  # Tamil
            return "ta"
    return None


def extract_ppt(
        ppt_file: Path,
        output_base: Path,
        project_format: str,
        chapter_regex: str,
        meta_prefix: str,
        meta_dir_sa: Path,
        meta_dir_ta: Path,
):
    """Scan PPT, auto-detect Sanskrit/Tamil shapes, and export text into lang folders."""
    # ✅ Determine chapter number
    chap_num = 1 if project_format == "standalone" else extract_chapter_num(ppt_file.stem, chapter_regex)
    if chap_num is None:
        print(f"⚠️ Could not parse chapter number from {ppt_file.name}")
        return
    chapter_id = f"chapter{chap_num:03d}"

    # ✅ Metadata lookup for Sanskrit
    meta_data_sa = None
    if meta_dir_sa and meta_dir_sa.exists():
        meta_file = meta_dir_sa / f"{meta_prefix}{chap_num:02d}.txt"
        if meta_file.exists():
            meta_data_sa = get_metadata(meta_file)
    # ✅ Metadata lookup for Tamil
    meta_data_ta = None
    if meta_dir_ta and meta_dir_ta.exists():
        meta_file = meta_dir_ta / f"{meta_prefix}{chap_num:02d}.txt"
        if meta_file.exists():
            meta_data_ta = get_metadata(meta_file)

    # ✅ Prepare output files
    lang_files = {}
    for lang, meta_data in [("sa", meta_data_sa), ("ta", meta_data_ta)]:
        out_dir = output_base / lang
        out_dir.mkdir(parents=True, exist_ok=True)
        out_file = out_dir / f"{chapter_id}.txt"
        lang_files[lang] = open(out_file, "w", encoding="utf-8")
        if meta_data:
            for meta_line in meta_data:
                lang_files[lang].write(meta_line + "\n")

    # ✅ Scan slides
    prs = Presentation(ppt_file)
    for i in range(1, len(prs.slides)-1):  # skip first slide
        slide = prs.slides[i]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            lang = detect_language(text)
            if lang in lang_files:
                lang_files[lang].write(text + "\n")
            else:
                print(f"⚠️ Slide {i}: Unknown language → {text[:30]}...")

    for f in lang_files.values():
        f.close()
    print(f"📖 Extracted {ppt_file.name} → {output_base}")


def extract_text(project_id: str):
    """Step 2: Extract text for all volumes as per orchestrator.plist."""
    project_dir = get_project_path(project_id)
    orch_plist = project_dir / "orchestrator.plist"
    with open(orch_plist, "rb") as f:
        config = plistlib.load(f)

    project_format = config.get("format", "volume")
    staging_root = Path(config.get("staging_root", ""))
    if not staging_root:
        raise ValueError("❌ staging_root not defined in orchestrator.plist")

    step_cfg = config.get("text_extraction", {})
    volumes_cfg = step_cfg.get("volumes", {})
    step_regex = step_cfg.get("chapter_regex", r"(\d+)")

    for vol_id, vol_cfg in volumes_cfg.items():
        if not vol_cfg.get("orchestrate", False):
            continue
        ppts_path = Path(vol_cfg.get("ppts_dir", ""))
        if not ppts_path.exists():
            print(f"⚠️ PPTs not found for {vol_id}: {ppts_path}")
            continue
        output_dir = staging_root / "content" / vol_id
        output_dir.mkdir(parents=True, exist_ok=True)
        files_to_process = [ppts_path] if ppts_path.is_file() else sorted(
            f for f in ppts_path.glob("*") if f.suffix.lower() in [".pptm", ".pptx", ".key"] and not f.name.startswith(('.', '~'))
        )
        if not files_to_process:
            continue
        meta_prefix = vol_cfg.get("meta_prefix", "")
        meta_dir_sa = Path(vol_cfg.get("meta_dir_sa", "")) if vol_cfg.get("meta_dir_sa") else None
        meta_dir_ta = Path(vol_cfg.get("meta_dir_ta", "")) if vol_cfg.get("meta_dir_ta") else None
        for ppt_file in files_to_process:
            if ppt_file.suffix.lower() == ".key":
                ppt_file = convert_key_to_pptx(ppt_file)
            regex = vol_cfg.get("chapter_regex", step_regex)
            extract_ppt(ppt_file, output_dir, project_format, regex, meta_prefix, meta_dir_sa, meta_dir_ta)


if __name__ == "__main__":
    extract_text("skp_srimad_bhagavatam")
